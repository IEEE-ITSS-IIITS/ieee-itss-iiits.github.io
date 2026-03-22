import collections.abc
import pptx
import os
import sys

def extract_pptx(filepath, output_dir):
    try:
        prs = pptx.Presentation(filepath)
    except Exception as e:
        print(f"Error reading pptx: {e}")
        sys.exit(1)

    os.makedirs(output_dir, exist_ok=True)
    
    with open("pptx_content.txt", "w", encoding="utf-8") as f:
        for i, slide in enumerate(prs.slides):
            f.write(f"--- Slide {i+1} ---\n")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    f.write(shape.text + "\n")
                if hasattr(shape, "image"):
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        image_filename = f"slide_{i+1}_image_{shape.shape_id}.{image.ext}"
                        with open(os.path.join(output_dir, image_filename), "wb") as img_file:
                            img_file.write(image_bytes)
                        f.write(f"[Image extracted: {image_filename}]\n")
                    except Exception as e:
                        f.write(f"[Error extracting image: {e}]\n")

if __name__ == "__main__":
    pptx_path = r"c:\Users\HARIHARAN\Desktop\ITSS\IEEE ITSS.pptx"
    out_dir = r"c:\Users\HARIHARAN\Desktop\ITSS\eventin-clone\public\assets"
    extract_pptx(pptx_path, out_dir)
    print("Extraction complete. Text saved to pptx_content.txt and images to eventin-clone/public/assets.")
